import random
from pptx import Presentation
from pptx.util import Inches

presentation = Presentation()

methods_to_describe = [
    ('random', 'Генерирует случайное число от 0 до 1.'),
    ('randint', 'Возвращает случайное целое число в заданном диапазоне.'),
    ('choice', 'Возвращает случайный элемент из последовательности.'),
    ('shuffle', 'Перемешивает элементы в последовательности.'),
    ('uniform', 'Возвращает случайное число в заданном диапазоне.')
]

methods_to_describe.extend([
    ('randrange', 'Возвращает случайное число в заданном диапазоне с заданным шагом.'),
    ('sample', 'Возвращает новый список, содержащий уникальные элементы из последовательности.'),
    ('seed', 'Устанавливает стартовое значение для генератора случайных чисел.'),
    ('getstate', 'Возвращает состояние генератора случайных чисел.')
])

for method, description in methods_to_describe:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5]) 
    title = slide.shapes.title
    content = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    content_text_frame = content.text_frame
    
    title.text = method
    
    method_description = f"Метод: {method}\n\nОписание:\n{description}"
    
    p = content_text_frame.add_paragraph()
    p.text = method_description
    p.space_after = 14

presentation.save('test.pptx')
